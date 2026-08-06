#!/usr/bin/env python3
"""
기획서 청크 추출 도구
PPTX, PDF를 지정된 크기로 청크로 쪼개고 출처 태그를 붙인다.
"""

import argparse
import json
import os
from pathlib import Path
from typing import Optional
import re
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("python-pptx가 필요합니다: pip install python-pptx")
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    import pdfplumber
except ImportError:
    print("pdfplumber가 필요합니다: pip install pdfplumber")
    pdfplumber = None


def get_slide_title(slide) -> str:
    """슬라이드의 제목을 추출한다."""
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip()
    return ""


def has_table_in_slide(slide) -> bool:
    """슬라이드에 명세 표가 있는지 확인한다."""
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            # 2x2 이상의 테이블만 명세 표로 판단
            if len(table.rows) >= 2 and len(table.columns) >= 2:
                return True
    return False


def extract_pptx(file_path: str) -> list[dict]:
    """PPTX에서 텍스트와 발표자 노트를 추출한다."""
    if not Presentation:
        raise ImportError("python-pptx를 설치해주세요: pip install python-pptx")

    prs = Presentation(file_path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_content = []
        has_image = False
        title = get_slide_title(slide)
        has_table = has_table_in_slide(slide)

        # 슬라이드 콘텐츠 추출
        for shape in slide.shapes:
            # 이미지 감지
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True

            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # 테이블 처리
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                slide_content.append("\n".join(rows))

        # 발표자 노트 추출
        notes = ""
        if slide.has_notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes = notes_text_frame.text.strip()

        # needs_vision 판정: (텍스트 30자 미만) OR (이미지 있음)
        content_text = "\n".join(slide_content)
        needs_vision = len(content_text.strip()) < 30 or has_image

        slides_data.append({
            "slide_num": slide_idx,
            "content": content_text,
            "notes": notes,
            "source": f"슬라이드 {slide_idx}",
            "title": title,
            "has_image": has_image,
            "has_table": has_table,
            "needs_vision": needs_vision
        })

    return slides_data


def extract_pdf(file_path: str) -> list[dict]:
    """PDF에서 텍스트를 추출한다."""
    if not pdfplumber:
        raise ImportError("pdfplumber를 설치해주세요: pip install pdfplumber")

    pages_data = []

    with pdfplumber.open(file_path) as pdf:
        for page_idx, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""

            # 테이블 추출
            tables = page.extract_tables()
            table_text = ""
            if tables:
                for table in tables:
                    rows = []
                    for row in table:
                        cells = [str(cell or "").strip() for cell in row]
                        rows.append(" | ".join(cells))
                    table_text += "\n" + "\n".join(rows)

            content = text + table_text

            # 이미지 감지 (pdfplumber)
            has_image = bool(page.images) if hasattr(page, 'images') else False

            # needs_vision 판정: (텍스트 30자 미만) OR (이미지 있음)
            needs_vision = len(content.strip()) < 30 or has_image

            pages_data.append({
                "page_num": page_idx,
                "content": content.strip(),
                "source": f"페이지 {page_idx}",
                "has_image": has_image,
                "needs_vision": needs_vision
            })

    return pages_data


def parse_pages(pages_str: str) -> list[int]:
    """'4-11,13-23' 형태의 페이지 범위 문자열을 숫자 리스트로 변환한다."""
    pages = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            try:
                start, end = part.split("-")
                start, end = int(start.strip()), int(end.strip())
                pages.extend(range(start, end + 1))
            except ValueError:
                print(f"[ERROR] 잘못된 페이지 범위 형식: {part}")
                return []
        else:
            try:
                pages.append(int(part))
            except ValueError:
                print(f"[ERROR] 잘못된 페이지 번호: {part}")
                return []
    return sorted(list(set(pages)))


def should_exclude(title: str, has_table: bool, is_pptx: bool) -> tuple[bool, Optional[str]]:
    """슬라이드/페이지를 제외할지 판단한다.

    Returns:
        (제외 여부, 제외 사유)
    """
    exclude_words = [
        "appendix", "개요", "목차", "표지", "cover", "index",
        "revision", "변경이력", "참고", "이력"
    ]

    # 제외어 체크
    title_lower = title.lower()
    for word in exclude_words:
        if word.lower() in title_lower:
            return True, f"제외어: {word}"

    # PPTX에서만 명세 표 체크
    if is_pptx and not has_table:
        return True, "명세 표 없음"

    return False, None


def filter_units(data: list[dict], is_pptx: bool) -> tuple[list[dict], list[dict]]:
    """필터 규칙을 적용하여 분석 대상과 제외 대상을 분리한다.

    Returns:
        (분석 대상 리스트, 제외 대상 리스트)
    """
    analyzed = []
    excluded = []

    for item in data:
        unit_num = item.get("slide_num") or item.get("page_num")
        title = item.get("title", "")
        has_table = item.get("has_table", False)

        should_exclude_flag, reason = should_exclude(title, has_table, is_pptx)

        if should_exclude_flag:
            excluded.append({
                "no": unit_num,
                "title": title,
                "reason": reason
            })
        else:
            analyzed.append(item)

    return analyzed, excluded


def add_source_tags(data: list[dict]) -> list[dict]:
    """모든 텍스트에 출처 태그를 붙인다."""
    for item in data:
        source = item["source"]
        content_lines = item["content"].split("\n")
        tagged_lines = [f"[{source}] {line}" if line.strip() else "" for line in content_lines]
        item["content"] = "\n".join(tagged_lines)

    return data


def chunk_content(data: list[dict], chunk_size: int = 8, overlap: int = 1) -> list[dict]:
    """콘텐츠를 청크로 쪼개고 중첩을 추가한다."""
    chunks = []
    chunk_id = 1

    # 모든 항목을 연결
    all_content = []
    for item in data:
        all_content.append(item)

    # 청크 생성
    for i in range(0, len(all_content), max(1, chunk_size - overlap)):
        chunk_items = all_content[i:i + chunk_size]

        if not chunk_items:
            break

        # 청크 메타데이터 생성
        chunk_text = "\n\n".join([item["content"] for item in chunk_items])
        sources = [item["source"] for item in chunk_items]
        needs_vision = any(item.get("needs_vision", False) for item in chunk_items)

        chunks.append({
            "id": f"C{chunk_id:03d}",
            "chunk_num": chunk_id,
            "content": chunk_text,
            "sources": sources,
            "source_range": f"{sources[0]} ~ {sources[-1]}" if sources else "",
            "needs_vision": needs_vision,
            "item_count": len(chunk_items)
        })

        chunk_id += 1

    return chunks


def compress_numbers(nums: list[int]) -> str:
    """연속된 숫자를 범위 형태로 축약한다. [4,5,6,7,11,13,14] -> '4-7, 11, 13-14'"""
    if not nums:
        return ""

    nums = sorted(set(nums))
    groups = []
    start = nums[0]
    end = nums[0]

    for num in nums[1:]:
        if num == end + 1:
            end = num
        else:
            if start == end:
                groups.append(str(start))
            else:
                groups.append(f"{start}-{end}")
            start = num
            end = num

    # 마지막 그룹 처리
    if start == end:
        groups.append(str(start))
    else:
        groups.append(f"{start}-{end}")

    return ", ".join(groups)


def print_dry_run_results(analyzed: list[dict], excluded: list[dict]):
    """dry-run 결과를 출력한다."""
    analyzed_nums = sorted(set([item.get("slide_num") or item.get("page_num") for item in analyzed]))
    compressed_range = compress_numbers(analyzed_nums)

    print(f"\n분석 대상 {len(analyzed_nums)}장: {compressed_range}")

    if excluded:
        print(f"제외 {len(excluded)}장:")
        # 연속된 번호 + 같은 제목끼리 그룹화
        current_group = []
        for exc in sorted(excluded, key=lambda x: x["no"]):
            if (not current_group or
                (exc["no"] == current_group[-1]["no"] + 1 and
                 exc["title"] == current_group[-1]["title"] and
                 exc["reason"] == current_group[-1]["reason"])):
                current_group.append(exc)
            else:
                # 그룹 출력
                if len(current_group) == 1:
                    item = current_group[0]
                    print(f"  {item['no']:3d}     {item['title']} ({item['reason']})")
                else:
                    nos = f"{current_group[0]['no']}-{current_group[-1]['no']}"
                    title = current_group[0]['title']
                    reason = current_group[0]['reason']
                    print(f"  {nos:3s}     {title} ({reason})")
                current_group = [exc]

        # 마지막 그룹 처리
        if current_group:
            if len(current_group) == 1:
                item = current_group[0]
                print(f"  {item['no']:3d}     {item['title']} ({item['reason']})")
            else:
                nos = f"{current_group[0]['no']}-{current_group[-1]['no']}"
                title = current_group[0]['title']
                reason = current_group[0]['reason']
                print(f"  {nos:3s}     {title} ({reason})")
    else:
        print("제외된 항목 없음")


def sanitize_filename(filename: str, max_length: int = 40) -> str:
    """파일명을 정규화한다: 특수문자 제거, 길이 제한."""
    # 길이 제한
    if len(filename) > max_length:
        filename = filename[:max_length]

    # 특수문자를 언더스코어로 치환
    special_chars = r'/ \ : * ? " < > |'
    for char in special_chars:
        filename = filename.replace(char, "_")

    # 연속된 언더스코어 정리
    while "__" in filename:
        filename = filename.replace("__", "_")

    # 앞뒤 공백 제거
    filename = filename.strip()

    return filename


def save_chunks(chunks: list[dict], output_dir: str, analyzed_units: list[int], excluded_units: list[dict], doc_name: str):
    """청크를 파일로 저장하고 manifest.json을 생성한다."""
    sanitized_name = sanitize_filename(doc_name)
    output_path = Path(output_dir) / sanitized_name
    chunks_dir = output_path / "chunks"
    chunks_dir.mkdir(parents=True, exist_ok=True)

    # vision이 필요한 페이지/슬라이드 수집 (확장된 기준 적용)
    vision_required_units = []

    manifest = {
        "total_chunks": len(chunks),
        "created_at": datetime.now().isoformat(),
        "analyzed_units": analyzed_units,
        "excluded_units": excluded_units,
        "chunks": [],
        "vision_required_units": vision_required_units
    }

    for chunk in chunks:
        chunk_file = chunks_dir / f"{chunk['id']}.md"

        # 마크다운 헤더 생성
        header = f"# {chunk['id']}\n\n"
        header += f"**출처**: {chunk['source_range']}\n"
        header += f"**항목 수**: {chunk['item_count']}\n"

        if chunk["needs_vision"]:
            header += "⚠️ **vision 필요**: 이미지 기반 내용 포함\n"

        header += "\n---\n\n"

        # 파일 저장
        with open(chunk_file, "w", encoding="utf-8") as f:
            f.write(header + chunk["content"])

        # manifest에 추가
        manifest["chunks"].append({
            "id": chunk["id"],
            "file": f"chunks/{chunk['id']}.md",
            "sources": chunk["sources"],
            "needs_vision": chunk["needs_vision"]
        })

    # vision_required_units 구성 (needs_vision=true인 모든 출처)
    for chunk in chunks:
        if chunk["needs_vision"]:
            for source in chunk["sources"]:
                if source not in vision_required_units:
                    vision_required_units.append(source)

    # manifest.json 저장
    manifest_file = output_path / "manifest.json"
    with open(manifest_file, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    print(f"[OK] {len(chunks)}개 청크를 {chunks_dir}에 저장했습니다")
    print(f"[OK] vision 필요 페이지/슬라이드: {len(vision_required_units)}개 ({', '.join(vision_required_units)})")
    print(f"[OK] manifest.json을 {manifest_file}에 저장했습니다")
    print(f"\n[PATH] 저장 경로: {output_path.absolute()}")


def main():
    parser = argparse.ArgumentParser(description="기획서를 청크로 쪼개고 출처 태그를 붙인다")
    parser.add_argument("file", help="PPTX 또는 PDF 파일 경로")
    parser.add_argument(
        "-o", "--output",
        default="output",
        help="출력 디렉토리 (기본값: output)"
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=8,
        help="한 청크의 항목 수 (기본값: 8)"
    )
    parser.add_argument(
        "--overlap",
        type=int,
        default=1,
        help="청크 중첩 항목 수 (기본값: 1)"
    )
    parser.add_argument(
        "--pages",
        type=str,
        default=None,
        help="분석할 페이지/슬라이드 범위 (예: 4-11,13-23)"
    )
    parser.add_argument(
        "--no-filter",
        action="store_true",
        help="자동 필터를 비활성화하고 전체를 분석 대상으로 한다"
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="청크 파일을 만들지 않고 분류 결과만 출력한다"
    )

    args = parser.parse_args()

    # 파일 존재 확인
    if not os.path.exists(args.file):
        print(f"[ERROR] 파일을 찾을 수 없습니다: {args.file}")
        return

    file_path = Path(args.file)
    file_ext = file_path.suffix.lower()
    doc_name = file_path.stem  # 확장자 제외한 파일명
    is_pptx = file_ext == ".pptx"

    # 파일 타입에 따라 추출
    if is_pptx:
        print(f"[PPTX] PPTX 파일을 처리 중입니다: {args.file}")
        data = extract_pptx(args.file)
    elif file_ext == ".pdf":
        print(f"[PDF] PDF 파일을 처리 중입니다: {args.file}")
        data = extract_pdf(args.file)
    else:
        print(f"[ERROR] 지원하지 않는 파일 형식: {file_ext} (PPTX, PDF만 지원)")
        return

    print(f"[OK] {len(data)}개 항목을 추출했습니다")

    # 필터링 처리
    excluded_units = []
    analyzed_units = []

    if args.pages:
        # --pages 옵션: 지정된 범위만 분석
        pages = parse_pages(args.pages)
        if not pages:
            return
        data = [item for item in data if (item.get("slide_num") or item.get("page_num")) in pages]
        analyzed_units = sorted(pages)
        print(f"[OK] 지정된 범위로 필터링: {len(data)}개 항목")
    elif args.no_filter:
        # --no-filter: 전체 분석
        analyzed_units = sorted([item.get("slide_num") or item.get("page_num") for item in data])
        print(f"[OK] 필터링 비활성화: 전체 {len(data)}개 항목 분석")
    else:
        # 자동 필터 적용
        data, excluded_units = filter_units(data, is_pptx)
        analyzed_units = sorted([item.get("slide_num") or item.get("page_num") for item in data])
        print(f"[OK] 자동 필터 적용: {len(data)}개 분석 대상, {len(excluded_units)}개 제외")

    # dry-run: 청크를 만들지 않고 분류 결과만 출력
    if args.dry_run:
        print_dry_run_results(data, excluded_units)
        return

    # 출처 태그 추가
    data = add_source_tags(data)

    # 청크로 쪼개기
    chunks = chunk_content(data, chunk_size=args.chunk_size, overlap=args.overlap)
    print(f"✓ {len(chunks)}개 청크로 분할했습니다")

    # 저장
    save_chunks(chunks, args.output, analyzed_units, excluded_units, doc_name)


if __name__ == "__main__":
    main()
